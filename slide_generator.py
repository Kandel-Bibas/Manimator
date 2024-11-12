import time
from dotenv import load_dotenv
import random

import subprocess
import re
import os
from logging import *
from prompts import *
import google.generativeai as genai
import openai
import fitz
from prompts import *


from pptx.util import Pt
from pptx.util import Inches
from pptx import Presentation





load_dotenv()
openai.api_key = os.getenv("OPENAI_API_KEY")
google_api_key = os.getenv('GOOGLE_API_KEY')


#Read contents within a PDF
def read_pdf(file_path):
  text = ""
  with fitz.open(file_path) as pdf:
      for page_num in range(pdf.page_count):
          page = pdf[page_num]
          text += page.get_text("text") + "\n" 
  return text


#Parses the response in JSON Format
def parse_json(response):
  text = response.text
  text = re.sub('([\t\n]+|\s{2,})', '', text)

  text = re.sub('(\},|\",|"|[{}[[\]]|```json|```|slides: )', '', text)

  text = re.sub('(heading: |subtext: |example: )', '\n', text)

  content = re.split('\n', text)
  return content


def generate_json(pdf_path):
    pdf_content = read_pdf(pdf_path)
    model = GetModel()
    prompt = MAKE_JSON_PROMPT

    response = model.generate_content([prompt, pdf_content])

    return parse_json(response)

    


#Extract Python Code and the Class for Manim
def extract_code_blocks(text):
    pattern = r"```python(.*?)```"
    matches = re.findall(pattern, text, re.DOTALL)

    class_name = r"class (.*?)\("
    class_match = re.findall(class_name, text, re.DOTALL)

    return matches, class_match


#Write into a Python file the manim script
def create_python_file(response):
    code_blocks = extract_code_blocks(response)[0]
    for block in code_blocks:
        code = block.strip()
    filename = f"manim_script"

    with open(f"{filename}.py", 'w') as file:
        file.write(code)
    return filename



#Sets the Gemini Model
def GetModel(api_key = google_api_key):
  genai.configure(api_key=api_key)
  generation_config = {
    "temperature": 1,
    "top_p": 0.9, 
    "top_k": 40,
    "response_mime_type": "text/plain",
  }
  model = genai.GenerativeModel(
    model_name="gemini-1.5-pro-002",
    generation_config=generation_config,
  )
  return model


def send_message_with_retries(chat, request, max_retries=3):
    retry_wait = 40 
    for attempt in range(max_retries):
        try:
            response = chat.send_message(request)
            return response.text
        except Exception as e:
            print(f"An error occurred: {e}")
            time.sleep(retry_wait)
            retry_wait *= 2 
    raise Exception("Max retries exceeded")





#Makes the Manim Animation, Tries to fix any issue with Gemini looking through the Frames.
def make_animation(topic):      

    prompt = ANIMATION_PROMPT_SYSTEM + ' ' + ANIMATION_PROMPT_USER.format(concept=topic) 
    model = GetModel(google_api_key)

    chat_session = model.start_chat(
        history=[]
    )

    attempted = 0
    completed = False
    next_prompt = prompt

    while (attempted < 5 and not completed):
            response = send_message_with_retries(chat_session, next_prompt)

            filename = create_python_file(response)
            class_name = extract_code_blocks(response)[1][0]

            command = f"manim -ql {filename}.py {class_name} --disable_caching"

            result = subprocess.run(command, shell=True, capture_output=True, text=True)

            if result.returncode == 0:
                completed = True
                print('Video Animated Successfully.')


            else:
                print(f'Failed Attempt {attempted} - Trying Again!')
                attempted += 1

                error_prompt = ERROR_PROMPT_SYSTEM.format(error=result.stderr) + ' ' + ERROR_PROMPT_USER
                next_prompt = "\n\n" + error_prompt

            if not completed and attempted ==5:
                print("Failed to generate a successful output even after 5 attempts.")


    all_videos_path = 'media/videos/manim_script/480p15/'
    file_video = all_videos_path + class_name + '.mp4'

    if os.path.exists(file_video):
      return file_video
    else:
        return 'content/Transparent_Square_Tiles_Texture.png'






def generate_pptx(pptx, refined_content):
  transparent_image = 'content/Transparent_Square_Tiles_Texture.png'

  for i in range(1, len(refined_content), 3):
    layout = pptx.slide_layouts[7]
    slide = pptx.slides.add_slide(layout).shapes
    title = slide.title
    title.text = refined_content[i]
    title.text_frame.paragraphs[0].font.size = Pt(36)

    body = slide.placeholders[2]
    subtext = body.text_frame.add_paragraph()
    subtext.text = refined_content[i + 1]
    subtext.font.size = Pt(18)

    
    if (i > 1):
        try:
            concept = refined_content[i + 1] + ' ' + refined_content[i + 2]
        except:
            concept = refined_content[i+1]
            
        video_file = make_animation(concept)
        video = slide.add_movie(video_file, Inches(6), Inches(2), Inches(5), Inches(2.8125), transparent_image, 'video/mp4')
        


  xml_slides = pptx.slides._sldIdLst
  slides = list(xml_slides)
  xml_slides.remove(slides[0])
  return pptx




def make_animation_and_slide(pdf_file_path):

    pdf = pdf_file_path
    file_name = pdf.split('.')[0] + '.pptx'

    refined_content = generate_json(pdf)
    print(len(refined_content))

    pptx = Presentation(f'content/template_3.pptx')


    pptx = generate_pptx(pptx, refined_content)

    pptx.save(file_name)

    print(file_name)
    return file_name

